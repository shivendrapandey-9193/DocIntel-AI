import os
import sys
import tempfile
import re
import logging
import html
from typing import List, Dict, Optional, Union, Any, Tuple
from pathlib import Path
from dotenv import load_dotenv
import requests
from bs4 import BeautifulSoup
from pypdf import PdfReader
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
import pandas as pd
from PIL import Image, ImageDraw, ImageFont, ImageEnhance
import pytesseract
import numpy as np
import fitz  # PyMuPDF
import networkx as nx
import matplotlib.pyplot as plt
from io import BytesIO
import json
import base64
import asyncio
import aiohttp
from cachetools import TTLCache
import smtplib
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart

# Optional imports
try:
    from odf import text, teletype
    from odf.opendocument import load as load_odt
except ImportError:
    text = None
    teletype = None
    load_odt = None

try:
    import docx2txt
except ImportError:
    docx2txt = None

try:
    import rarfile
except ImportError:
    rarfile = None

try:
    import py7zr
except ImportError:
    py7zr = None

try:
    import yake
except ImportError:
    yake = None

import subprocess
from collections import Counter

# LangChain imports
try:
    from langchain_core.documents import Document as LangchainDocument
    from langchain_text_splitters import RecursiveCharacterTextSplitter
    from langchain_community.vectorstores import FAISS
    from langchain_core.prompts import ChatPromptTemplate
    from langchain_core.runnables import RunnablePassthrough
    from langchain_core.output_parsers import StrOutputParser
    from langchain_community.document_loaders import WebBaseLoader
    from langchain_together import Together
    from langchain_community.embeddings import HuggingFaceEmbeddings
    from langchain.memory import ConversationBufferMemory
except ImportError:
    pass

# Streamlit imports
import streamlit as st
try:
    from streamlit_option_menu import option_menu
except ImportError:
    option_menu = None
try:
    import streamlit_lottie as st_lottie
except ImportError:
    st_lottie = None

# SerpAPI import
try:
    from serpapi import GoogleSearch
except ImportError:
    GoogleSearch = None

# For archive handling
import zipfile
import tarfile

# For advanced conversions
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from docx.shared import Pt, Inches
from docx.enum.text import WD_PARAGRAPH_ALIGNMENT

# Wordcloud for keyword visualization
try:
    from wordcloud import WordCloud
except ImportError:
    WordCloud = None

# Suppress warnings
import warnings
warnings.filterwarnings("ignore", category=FutureWarning)
warnings.filterwarnings("ignore", category=DeprecationWarning)
warnings.filterwarnings("ignore", category=UserWarning, module="langchain_core")
warnings.filterwarnings("ignore", category=UserWarning, module="streamlit")

# Custom logging filter to suppress answer content
class NoAnswerFilter(logging.Filter):
    def filter(self, record):
        return not any(keyword in record.msg.lower() for keyword in ["answer", "response", "content"])

# Logging configuration
LOG_FILE = "docintel.log"
logger = logging.getLogger("docintel")
logger.setLevel(logging.INFO)

if not logger.handlers:
    fh = logging.FileHandler(LOG_FILE)
    fh.setLevel(logging.INFO)
    sh = logging.StreamHandler(sys.stdout)
    sh.setLevel(logging.INFO)
    formatter = logging.Formatter(
        "[%(asctime)s] %(levelname)s in %(module)s: %(message)s"
    )
    fh.setFormatter(formatter)
    sh.setFormatter(formatter)
    sh.addFilter(NoAnswerFilter())  # Suppress answer content in logs
    logger.addHandler(fh)
    logger.addHandler(sh)

# Set Tesseract path
TESSERACT_PATH = os.getenv("TESSERACT_PATH", r"C:\Program Files\Tesseract-OCR\tesseract.exe")
try:
    pytesseract.pytesseract.tesseract_cmd = TESSERACT_PATH
    if not os.path.exists(TESSERACT_PATH):
        raise FileNotFoundError(f"Tesseract not found at {TESSERACT_PATH}")
    logger.info(f"Tesseract found at {TESSERACT_PATH}")
except Exception as e:
    logger.warning(f"Tesseract not found: {e}. OCR features disabled. Ensure Tesseract is installed and the path is correct.")
    pytesseract.pytesseract.tesseract_cmd = None

# Global spaCy model singleton
_nlp = None
def load_spacy_model(model_name: str = "en_core_web_md"):  # Upgraded to medium model for better accuracy
    global _nlp
    if _nlp is not None:
        logger.info(f"Using cached spaCy model {model_name}")
        return _nlp
    try:
        import spacy
        try:
            _nlp = spacy.load(model_name)
            logger.info(f"Loaded spaCy model {model_name}")
            return _nlp
        except OSError:
            logger.warning(f"spaCy model {model_name} not found, attempting to download...")
            subprocess.run([sys.executable, "-m", "spacy", "download", model_name], check=True)
            _nlp = spacy.load(model_name)
            logger.info(f"Successfully downloaded and loaded spaCy model {model_name}")
            return _nlp
    except Exception as ex:
        logger.error(f"Failed to load spaCy model {model_name}: {ex}")
    return None

# Initialize spaCy model once
nlp = load_spacy_model("en_core_web_md")

# Cache for web content and LLM responses
web_cache = TTLCache(maxsize=100, ttl=3600)  # Cache web content for 1 hour
llm_cache = TTLCache(maxsize=1000, ttl=1800)  # Cache LLM responses for 30 minutes

# Requests session with global User-Agent
if "USER_AGENT" not in os.environ:
    os.environ["USER_AGENT"] = "DocIntelAI/1.0 (Python; +https://github.com/docintel-ai)"
USER_AGENT = os.environ.get("USER_AGENT")
session = requests.Session()
session.headers.update({"User-Agent": USER_AGENT})

# Async web content fetching with improved image validation
async def fetch_url_content(url: str) -> Dict[str, Any]:
    if url in web_cache:
        logger.info(f"Returning cached content for {url}")
        return web_cache[url]
    # Skip LinkedIn URLs to avoid 999 errors
    if "linkedin.com" in url.lower():
        return {"text": "LinkedIn URLs are not supported due to access restrictions. Please provide an alternative source.", "links": [], "images": [], "videos": []}
    try:
        async with aiohttp.ClientSession(headers={"User-Agent": USER_AGENT}) as session:
            async with session.get(url, timeout=10) as response:
                if response.status != 200:
                    logger.error(f"Failed to fetch {url}: Status {response.status}")
                    return {"text": "", "links": [], "images": [], "videos": []}
                html_content = await response.text()
                soup = BeautifulSoup(html_content, 'html.parser')
                
                # Extract main content
                main_content = ""
                main_tags = soup.find_all(['p', 'article', 'div'], class_=re.compile('content|main|body', re.I))
                for tag in main_tags:
                    main_content += tag.get_text(strip=True) + "\n"
                if not main_content.strip():
                    main_content = soup.get_text(strip=True)
                
                # Extract links
                links = [a['href'] for a in soup.find_all('a', href=True) if a['href'].startswith(('http://', 'https://'))]
                
                # Extract images with validation
                images = []
                for img in soup.find_all('img', src=True):
                    src = img['src']
                    if src.startswith(('http://', 'https://')):
                        try:
                            async with session.get(src) as img_response:
                                if img_response.status == 200:
                                    img_data = await img_response.read()
                                    # Validate if it's a valid image
                                    try:
                                        Image.open(BytesIO(img_data))
                                        images.append({
                                            "bytes": BytesIO(img_data),
                                            "source": url,
                                            "format": src.split('.')[-1].lower(),
                                            "width": img.get('width', 'Unknown'),
                                            "height": img.get('height', 'Unknown')
                                        })
                                    except:
                                        logger.warning(f"Invalid image data from {src}")
                        except Exception as e:
                            logger.error(f"Error fetching image {src}: {e}")
                
                # Extract videos
                videos = []
                for video in soup.find_all(['video', 'source'], src=True):
                    src = video['src']
                    if src.startswith(('http://', 'https://')):
                        videos.append(src)
                
                result = {
                    "text": DocumentProcessor.clean_text(main_content),
                    "links": links[:10],  # Limit to 10 links
                    "images": images,
                    "videos": videos
                }
                web_cache[url] = result
                return result
    except Exception as e:
        logger.error(f"Error fetching URL {url}: {e}")
        return {"text": "", "links": [], "images": [], "videos": []}

def get_image_download_link(img_bytes: bytes, filename: str = "image.png", text: str = "Download") -> str:
    try:
        b64 = base64.b64encode(img_bytes).decode()
        return f'<a href="data:image/png;base64,{b64}" download="{html.escape(filename)}">{html.escape(text)}</a>'
    except Exception as e:
        logger.error(f"Error generating download link: {e}")
        return ""

def create_watermarked_image(image_bytes: bytes, text: str = "DocIntel AI") -> bytes:
    try:
        img = Image.open(BytesIO(image_bytes)).convert("RGBA")
        width, height = img.size
        overlay = Image.new("RGBA", img.size, (255, 255, 255, 0))
        draw = ImageDraw.Draw(overlay)
        font_size = max(18, int(min(width, height) * 0.04))
        try:
            font = ImageFont.truetype("arial.ttf", font_size)
        except:
            try:
                font = ImageFont.truetype("DejaVuSans.ttf", font_size)
            except:
                font = ImageFont.load_default()
        try:
            text_bbox = draw.textbbox((0, 0), text, font=font)
            text_width = text_bbox[2] - text_bbox[0]
            text_height = text_bbox[3] - text_bbox[1]
        except:
            text_width, text_height = draw.textsize(text, font=font)
        padding = int(font_size * 0.4)
        x = width - text_width - padding
        y = height - text_height - padding
        shadow_position = (x + 1, y + 1)
        draw.text(shadow_position, text, font=font, fill=(0, 0, 0, 120))
        draw.text((x, y), text, font=font, fill=(255, 255, 255, 180))
        watermarked = Image.alpha_composite(img, overlay).convert("RGB")
        buf = BytesIO()
        watermarked.save(buf, format="PNG")
        buf.seek(0)
        return buf.getvalue()
    except Exception as e:
        logger.error(f"Error creating watermarked image: {e}")
        return image_bytes

# Jabalpur Smart City Image Fetch
def fetch_jabalpur_image() -> Optional[BytesIO]:
    url = "https://jabalpursmartcity.in/"
    try:
        content = asyncio.run(fetch_url_content(url))
        if content["images"]:
            img_data = content["images"][0]["bytes"].getvalue()
            img = Image.open(BytesIO(img_data)).resize((300, 300))
            buf = BytesIO()
            img.save(buf, format="PNG")
            buf.seek(0)
            return buf
    except:
        logger.warning("No image found for Jabalpur Smart City.")
    return None

# Document Processing with improved OCR check and archive handling
class DocumentProcessor:
    @staticmethod
    def get_supported_formats() -> Dict[str, List[str]]:
        return {
            "📄 Documents": [".pdf", ".docx", ".txt", ".doc", ".rtf", ".odt", ".tex"],
            "📊 Spreadsheets": [".xlsx", ".csv", ".xls", ".tsv", ".ods", ".xlsm"],
            "🖼️ Images": [".jpg", ".jpeg", ".png", ".tiff", ".bmp"],
            "💻 Code": [".py", ".json", ".html", ".js", ".css", ".java"],
            "📦 Archives": [".zip", ".tar", ".gz", ".rar", ".7z"]
        }

    @staticmethod
    def extract_from_archive(file_path: str) -> List[str]:
        ext = os.path.splitext(file_path)[1].lower()
        temp_dir = tempfile.mkdtemp()
        extracted_files = []
        try:
            if ext == ".zip":
                with zipfile.ZipFile(file_path, "r") as zip_ref:
                    zip_ref.extractall(temp_dir)
            elif ext in [".tar", ".gz"]:
                with tarfile.open(file_path) as tar_ref:
                    tar_ref.extractall(temp_dir)
            elif ext == ".rar":
                if rarfile is None:
                    raise ImportError("rarfile not available")
                with rarfile.RarFile(file_path) as rar_ref:
                    rar_ref.extractall(temp_dir)
            elif ext == ".7z":
                if py7zr is None:
                    raise ImportError("py7zr not available")
                with py7zr.SevenZipFile(file_path, mode='r') as z:
                    z.extractall(temp_dir)
            for root, _, files in os.walk(temp_dir):
                for file in files:
                    extracted_path = os.path.join(root, file)
                    extracted_files.append(extracted_path)
        except Exception as e:
            logger.error(f"Error extracting archive {file_path}: {e}")
        return extracted_files

    @staticmethod
    def extract_images_from_pdf(pdf_path: str) -> List[Dict[str, Any]]:
        images = []
        seen_hashes = set()
        try:
            with fitz.open(pdf_path) as pdf:
                for page_num in range(len(pdf)):
                    page = pdf.load_page(page_num)
                    for img in page.get_images(full=True):
                        xref = img[0]
                        base_image = pdf.extract_image(xref)
                        image_bytes = base_image["image"]
                        image_hash = hash(image_bytes)
                        if image_hash not in seen_hashes:
                            # Validate image
                            try:
                                Image.open(BytesIO(image_bytes))
                                images.append({
                                    "bytes": BytesIO(image_bytes),
                                    "source": f"Page {page_num+1}",
                                    "format": base_image["ext"],
                                    "width": base_image["width"],
                                    "height": base_image["height"]
                                })
                                seen_hashes.add(image_hash)
                            except:
                                logger.warning(f"Invalid image on page {page_num+1}")
        except Exception as e:
            logger.error(f"Error extracting images from PDF: {e}")
        return images

    @staticmethod
    def extract_images_from_docx(docx_path: str) -> List[Dict[str, Any]]:
        images = []
        seen_hashes = set()
        try:
            doc = DocxDocument(docx_path)
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    img_data = rel.target_part.blob
                    image_hash = hash(img_data)
                    if image_hash not in seen_hashes:
                        try:
                            with Image.open(BytesIO(img_data)) as img:
                                images.append({
                                    "bytes": BytesIO(img_data),
                                    "source": "Document",
                                    "format": img.format.lower(),
                                    "width": img.width,
                                    "height": img.height
                                })
                                seen_hashes.add(image_hash)
                        except:
                            logger.warning(f"Invalid image in DOCX")
        except Exception as e:
            logger.error(f"Error extracting images from DOCX: {e}")
        return images

    @staticmethod
    def extract_images(file_path: str) -> List[Dict[str, Any]]:
        ext = os.path.splitext(file_path)[1].lower()
        if ext == ".pdf":
            return DocumentProcessor.extract_images_from_pdf(file_path)
        elif ext in [".docx", ".doc"]:
            return DocumentProcessor.extract_images_from_docx(file_path)
        elif ext in [".jpg", ".jpeg", ".png", ".tiff", ".bmp"]:
            try:
                with open(file_path, "rb") as f:
                    img_data = f.read()
                img = Image.open(BytesIO(img_data))
                return [{
                    "bytes": BytesIO(img_data),
                    "source": "Image",
                    "format": ext[1:].upper(),
                    "width": img.width,
                    "height": img.height
                }]
            except Exception as e:
                logger.error(f"Error processing image file {file_path}: {e}")
                return []
        return []

    @staticmethod
    def convert_to_text(file_path: str, target_format: str = "txt") -> Optional[Union[str, BytesIO]]:
        if not os.path.exists(file_path):
            logger.error(f"File not found: {file_path}")
            return None
        ext = os.path.splitext(file_path)[1].lower()
        if ext in [".zip", ".tar", ".gz", ".rar", ".7z"]:
            extracted = DocumentProcessor.extract_from_archive(file_path)
            text = ""
            for ex_file in extracted:
                text += DocumentProcessor.convert_to_text(ex_file, "txt") or ""
            return text if target_format == "txt" else DocumentProcessor._convert_text_to_format(text, target_format)
        try:
            text = ""
            if ext == ".pdf":
                with fitz.open(file_path) as pdf:
                    text = "\n".join([page.get_text() for page in pdf])
                    if not text.strip():
                        text = DocumentProcessor._pdf_ocr(file_path)
            elif ext in [".docx", ".doc"]:
                if docx2txt is None:
                    raise ImportError("docx2txt not available for DOCX processing")
                text = docx2txt.process(file_path)
            elif ext == ".rtf":
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    text = f.read()
            elif ext == ".odt":
                if load_odt is None:
                    raise ImportError("ODF tools not available for ODT processing")
                doc = load_odt(file_path)
                text = teletype.extractText(doc)
            elif ext == ".tex":
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    text = f.read()
                text = re.sub(r'%.*?\n', '', text)
                text = re.sub(r'\\[a-zA-Z]+{', '{', text)
            elif ext in [".xlsx", ".xls", ".xlsm", ".ods"]:
                wb = load_workbook(file_path, read_only=True)
                text = ""
                for sheet in wb.sheetnames:
                    text += f"Sheet: {sheet}\n"
                    for row in wb[sheet].iter_rows(values_only=True):
                        text += "\t".join(str(cell) for cell in row if cell is not None) + "\n"
            elif ext == ".pptx":
                prs = Presentation(file_path)
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
            elif ext in [".jpg", ".jpeg", ".png", ".tiff", ".bmp"]:
                text = DocumentProcessor._enhanced_ocr(file_path)
            elif ext == ".csv":
                try:
                    df = pd.read_csv(file_path, encoding="utf-8", encoding_errors="ignore")
                    text = df.to_string()
                except Exception as e:
                    logger.error(f"Error reading CSV: {e}")
                    return None
            elif ext == ".tsv":
                df = pd.read_csv(file_path, sep="\t", encoding="utf-8", errors="ignore")
                text = df.to_string()
            else:
                with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                    text = f.read()
            text = DocumentProcessor.clean_text(text)
            return DocumentProcessor._convert_text_to_format(text, target_format)
        except Exception as e:
            logger.error(f"Error converting file to {target_format}: {e}")
            return None

    @staticmethod
    def _convert_text_to_format(text: str, target_format: str) -> Optional[Union[str, BytesIO]]:
        if target_format == "txt":
            return text
        elif target_format == "json":
            return json.dumps({"content": text}, ensure_ascii=False)
        elif target_format == "md":
            # Basic formatting preservation
            md_text = re.sub(r'\n{2,}', '\n\n', text)
            md_text = re.sub(r'^(.+?)$', r'## \1', md_text, flags=re.M) if len(md_text.split('\n')) < 5 else md_text
            return md_text
        elif target_format == "html":
            html_text = re.sub(r'\n', '<br>', html.escape(text))
            html_text = f"<h2>Converted Content</h2><p>{html_text}</p>"
            return html_text
        elif target_format == "xml":
            return f'<?xml version="1.0" encoding="UTF-8"?>\n<document>\n<content>{html.escape(text)}</content>\n</document>'
        elif target_format == "docx":
            try:
                doc = DocxDocument()
                heading = doc.add_heading('Converted Content', level=1)
                heading.alignment = WD_PARAGRAPH_ALIGNMENT.CENTER
                lines = text.split('\n')
                for line in lines:
                    if line.strip():
                        p = doc.add_paragraph(line)
                        if re.match(r'^[A-Z ]{3,}$', line):
                            p.bold = True
                            p.style.font.size = Pt(14)
                buf = BytesIO()
                doc.save(buf)
                buf.seek(0)
                return buf
            except Exception as e:
                logger.error(f"Error converting to DOCX: {e}")
                return None
        elif target_format == "pdf":
            try:
                buf = BytesIO()
                doc = SimpleDocTemplate(buf, pagesize=letter)
                styles = getSampleStyleSheet()
                story = []
                lines = text.split('\n')
                for line in lines:
                    if re.match(r'^[A-Z ]{3,}$', line):
                        story.append(Paragraph(line, styles['Heading1']))
                    else:
                        story.append(Paragraph(line, styles['Normal']))
                    story.append(Spacer(1, 12))
                doc.build(story)
                buf.seek(0)
                return buf
            except Exception as e:
                logger.error(f"Error converting to PDF: {e}")
                return None
        elif target_format == "image_to_pdf":
            # Assume text is from image OCR, convert to PDF
            try:
                buf = BytesIO()
                c = canvas.Canvas(buf, pagesize=letter)
                c.drawString(100, 750, text)
                c.save()
                buf.seek(0)
                return buf
            except Exception as e:
                logger.error(f"Error converting image to PDF: {e}")
                return None
        elif target_format == "image_to_docx":
            # Assume text is from image OCR, convert to DOCX
            try:
                doc = DocxDocument()
                doc.add_paragraph(text)
                buf = BytesIO()
                doc.save(buf)
                buf.seek(0)
                return buf
            except Exception as e:
                logger.error(f"Error converting image to DOCX: {e}")
                return None
        logger.error(f"Unsupported target format: {target_format}")
        return None

    @staticmethod
    def _enhanced_ocr(image_path: str) -> str:
        try:
            if not pytesseract.pytesseract.tesseract_cmd:
                raise Exception(f"Tesseract not found at {TESSERACT_PATH}. Please install Tesseract and ensure the path is correct.")
            img = Image.open(image_path)
            img = img.convert('L')
            img = ImageEnhance.Contrast(img).enhance(2.0)
            img = ImageEnhance.Sharpness(img).enhance(2.0)
            img = Image.eval(img, lambda x: 0 if x < 128 else 255)
            custom_config = r'--oem 3 --psm 6 -l eng'  # Added multi-language support if needed, e.g., -l eng+hin
            text = pytesseract.image_to_string(img, config=custom_config)
            return text.strip() if text.strip() else "No text detected in image"
        except Exception as e:
            logger.error(f"OCR failed for image: {e}")
            return f"OCR unavailable: {str(e)}. Ensure Tesseract is installed at {TESSERACT_PATH} and pytesseract is properly configured."

    @staticmethod
    def _pdf_ocr(pdf_path: str) -> str:
        text = ""
        try:
            if not pytesseract.pytesseract.tesseract_cmd:
                raise Exception(f"Tesseract not found at {TESSERACT_PATH}. Please install Tesseract and ensure the path is correct.")
            with fitz.open(pdf_path) as pdf:
                for page_num in range(len(pdf)):
                    pix = pdf[page_num].get_pixmap(matrix=fitz.Matrix(2, 2))
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                        img.save(tmp.name)
                        page_text = DocumentProcessor._enhanced_ocr(tmp.name)
                        text += f"\nPage {page_num+1}:\n{page_text}\n"
                    os.remove(tmp.name)
        except Exception as e:
            logger.error(f"PDF OCR failed: {e}")
            return f"OCR unavailable: {str(e)}. Ensure Tesseract is installed at {TESSERACT_PATH} and pytesseract is properly configured."
        return text

    @staticmethod
    def clean_text(text: str) -> str:
        text = re.sub(r'[^\x20-\x7E\n\t\r]', '', text)
        text = re.sub(r'[ \t]+', ' ', text)
        text = re.sub(r'\n\s+\n', '\n\n', text)
        return text.strip()

# Document Loader
class DocumentLoader:
    @staticmethod
    async def load_documents(file_paths: List[str]) -> List[LangchainDocument]:
        documents = []
        for file_path in file_paths:
            try:
                if file_path.startswith(('http://', 'https://')):
                    web_content = await fetch_url_content(file_path)
                    if web_content["text"].strip():
                        documents.append(LangchainDocument(
                            page_content=web_content["text"],
                            metadata={"source": file_path, "links": web_content["links"], "images": web_content["images"], "videos": web_content["videos"]}
                        ))
                    else:
                        logger.warning(f"No content extracted from URL: {file_path}")
                    continue
                ext = os.path.splitext(file_path)[1].lower()
                if ext in [".zip", ".tar", ".gz", ".rar", ".7z"]:
                    extracted = DocumentProcessor.extract_from_archive(file_path)
                    for ex_file in extracted:
                        text = DocumentProcessor.convert_to_text(ex_file)
                        if text and isinstance(text, str) and text.strip():
                            metadata = {"source": os.path.basename(ex_file)}
                            documents.append(LangchainDocument(page_content=text, metadata=metadata))
                else:
                    text = DocumentProcessor.convert_to_text(file_path)
                    if text and isinstance(text, str) and text.strip():
                        metadata = {"source": os.path.basename(file_path)}
                        documents.append(LangchainDocument(page_content=text, metadata=metadata))
                    else:
                        logger.error(f"No content extracted from {file_path}")
            except Exception as e:
                logger.error(f"Error loading document {file_path}: {e}")
        return documents

# NER Analyzer with Chart.js visualization
class NERAnalyzer:
    @staticmethod
    def analyze_documents(documents: List[LangchainDocument]) -> Dict[str, Any]:
        if nlp is None:
            logger.error("spaCy model not loaded")
            return {"entities": [], "graph": nx.Graph(), "raw_text": ""}
        combined_text = "\n\n".join([doc.page_content for doc in documents])
        doc = nlp(combined_text)
        entities = []
        seen_entities = set()
        for ent in doc.ents:
            entity_key = (ent.text, ent.label_)
            if entity_key not in seen_entities:
                entities.append({
                    "text": ent.text,
                    "label": ent.label_,
                    "start": ent.start_char,
                    "end": ent.end_char
                })
                seen_entities.add(entity_key)
        graph = NERAnalyzer._build_entity_graph(doc)
        return {
            "entities": entities,
            "graph": graph,
            "raw_text": combined_text
        }

    @staticmethod
    def _build_entity_graph(doc) -> nx.Graph:
        graph = nx.Graph()
        for ent in doc.ents:
            graph.add_node(ent.text, type=ent.label_)
        for sent in doc.sents:
            sent_ents = list(set([ent.text for ent in sent.ents]))
            for i, ent1 in enumerate(sent_ents):
                for ent2 in sent_ents[i+1:]:
                    if graph.has_edge(ent1, ent2):
                        graph[ent1][ent2]["weight"] += 1
                    else:
                        graph.add_edge(ent1, ent2, weight=1)
        return graph

    @staticmethod
    def visualize_entities(entities: List[Dict]) -> str:
        if not entities:
            return "<p style='color: #333; font-size: 1.2rem;'>No entities found</p>"
        entity_types = {}
        for entity in entities:
            if entity["label"] not in entity_types:
                entity_types[entity["label"]] = []
            entity_types[entity["label"]].append(entity["text"])
        color_map = {
            "PERSON": "#FFD700",
            "ORG": "#87CEEB",
            "GPE": "#98FB98",
            "DATE": "#FFA07A",
            "NORP": "#DDA0DD",
            "CARDINAL": "#FFB6C1",
            "PRODUCT": "#FF69B4",
            "PERCENT": "#FF4500",
            "ORDINAL": "#ADFF2F",
            "OTHER": "#CCCCCC"
        }
        html_content = "<div class='entity-container'>"
        html_content += "<h2 style='color: #333; margin-bottom: 20px;'>🔍 Named Entities</h2>"
        html_content += "<table class='entity-table'>"
        html_content += "<tr><th style='background-color: #4facfe;'>Entity Type</th><th style='background-color: #4facfe;'>Entity Text</th></tr>"
        for label in sorted(entity_types.keys()):
            color = color_map.get(label, color_map["OTHER"])
            for text in sorted(set(entity_types[label])):
                html_content += f"<tr style='background-color: {color};'><td>{label.upper()}</td><td>{html.escape(text)}</td></tr>"
        html_content += "</table></div>"
        # Add Chart.js bar chart for entity frequencies
        freq = Counter([ent["label"] for ent in entities])
        labels = list(freq.keys())
        data = list(freq.values())
        colors = [color_map.get(label, "#CCCCCC") for label in labels]
        chart_html = """
        <div style="margin-top: 30px;">
            <h3>📊 Entity Type Frequencies</h3>
            <canvas id="entityChart" width="400" height="200"></canvas>
            <script src="https://cdn.jsdelivr.net/npm/chart.js"></script>
            <script>
                var ctx = document.getElementById('entityChart').getContext('2d');
                var myChart = new Chart(ctx, {
                    type: 'bar',
                    data: {
                        labels: %s,
                        datasets: [{
                            label: 'Frequency',
                            data: %s,
                            backgroundColor: %s
                        }]
                    },
                    options: {
                        scales: {
                            y: { beginAtZero: true }
                        }
                    }
                });
            </script>
        </div>
        """ % (json.dumps(labels), json.dumps(data), json.dumps(colors))
        html_content += chart_html
        return html_content

    @staticmethod
    def visualize_graph(graph: nx.Graph) -> Optional[BytesIO]:
        if not graph.nodes():
            return None
        plt.figure(figsize=(16, 12))
        node_colors = []
        for node in graph.nodes():
            node_type = graph.nodes[node].get("type", "OTHER")
            colors = {
                "PERSON": "#FFD700",
                "ORG": "#87CEEB",
                "GPE": "#98FB98",
                "DATE": "#FFA07A",
                "NORP": "#DDA0DD",
                "CARDINAL": "#FFB6C1",
                "PRODUCT": "#FF69B4",
                "PERCENT": "#FF4500",
                "ORDINAL": "#ADFF2F"
            }
            node_colors.append(colors.get(node_type, "#CCCCCC"))
        pos = nx.spring_layout(graph, k=1.2, iterations=50, scale=2)
        edge_weights = [graph[u][v]["weight"] for u, v in graph.edges()]
        max_weight = max(edge_weights) if edge_weights else 1
        edge_widths = [5 * (w / max_weight) for w in edge_weights] if edge_weights else [1]
        nx.draw_networkx_nodes(graph, pos, node_size=1500, node_color=node_colors, alpha=0.85)
        if edge_widths:
            nx.draw_networkx_edges(graph, pos, width=edge_widths, alpha=0.7, edge_color="#555555")
        labels = {n: n for n in graph.nodes()}
        nx.draw_networkx_labels(graph, pos, labels, font_size=12, font_weight="bold")
        edge_labels = nx.get_edge_attributes(graph, 'weight')
        if edge_labels:
            nx.draw_networkx_edge_labels(graph, pos, edge_labels=edge_labels, font_size=10)
        # Add legend
        from matplotlib.patches import Patch
        legend_elements = [
            Patch(facecolor="#FFD700", label="Person"),
            Patch(facecolor="#87CEEB", label="Organization"),
            Patch(facecolor="#98FB98", label="Location"),
            Patch(facecolor="#FFA07A", label="Date"),
            Patch(facecolor="#DDA0DD", label="Group"),
            Patch(facecolor="#FFB6C1", label="Number"),
            Patch(facecolor="#FF69B4", label="Product"),
            Patch(facecolor="#FF4500", label="Percent"),
            Patch(facecolor="#ADFF2F", label="Ordinal")
        ]
        plt.legend(handles=legend_elements, loc="upper right", fontsize=10, title="Entity Types")
        plt.title("Entity Relationship Graph", fontsize=18, pad=20)
        plt.axis("off")
        buf = BytesIO()
        plt.savefig(buf, format="png", bbox_inches="tight", dpi=200)
        plt.close()
        buf.seek(0)
        return buf

# Keyword Analyzer with WordCloud visualization
class KeywordAnalyzer:
    @staticmethod
    def extract_keywords(text: str) -> List[Tuple[str, float]]:
        if yake is None:
            logger.warning("YAKE not available for keyword extraction")
            return []
        kw_extractor = yake.KeywordExtractor(lan="en", n=3, dedupLim=0.3, top=20, features=None)
        keywords = kw_extractor.extract_keywords(text)
        return keywords

    @staticmethod
    def visualize_keywords(keywords: List[Tuple[str, float]]) -> Optional[BytesIO]:
        if not keywords or WordCloud is None:
            logger.warning("WordCloud not available or no keywords")
            return None
        word_dict = {kw: 1 - score for kw, score in keywords}  # Invert score for importance
        wc = WordCloud(width=800, height=400, background_color="white", colormap="plasma", max_words=50).generate_from_frequencies(word_dict)
        buf = BytesIO()
        plt.figure(figsize=(10, 5))
        plt.imshow(wc, interpolation="bilinear")
        plt.axis("off")
        plt.savefig(buf, format="png", bbox_inches="tight")
        plt.close()
        buf.seek(0)
        return buf

# RAG Pipeline with SerpAPI integration and improved prompts
class RAGPipeline:
    def __init__(self):
        self.vectorstore = None
        self.llm = None
        self.retriever = None
        self.embeddings = None
        self.memory = ConversationBufferMemory(return_messages=True)
        try:
            self.embeddings = HuggingFaceEmbeddings(model_name="sentence-transformers/all-mpnet-base-v2")  # Upgraded embedding model for better accuracy
        except Exception as e:
            logger.error(f"Error initializing embeddings: {e}")

    def create_vectorstore(self, documents: List[LangchainDocument]) -> None:
        try:
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=1500,  # Reduced chunk size for more granular retrieval
                chunk_overlap=300
            )
            splits = text_splitter.split_documents(documents)
            if not splits:
                raise ValueError("No document chunks created")
            self.vectorstore = FAISS.from_documents(splits, self.embeddings)
            self.retriever = self.vectorstore.as_retriever(
                search_type="similarity_score_threshold",
                search_kwargs={"k": 12, "score_threshold": 0.7}  # Added threshold for better accuracy
            )
        except Exception as e:
            logger.error(f"Error creating vectorstore: {e}")
            raise

    def set_llm(self) -> None:
        try:
            api_key = os.getenv("TOGETHER_API_KEY")
            if not api_key:
                raise ValueError("TOGETHER_API_KEY environment variable not set")
            self.llm = Together(
                model="mistralai/Mixtral-8x7B-Instruct-v0.1",
                temperature=0.5,  # Lowered temperature for more accurate responses
                max_tokens=2048,  # Increased tokens for detailed answers
                top_k=40,
                top_p=0.9,
                together_api_key=api_key
            )
        except Exception as e:
            logger.error(f"Error initializing LLM: {e}")
            raise

    def search_with_serpapi(self, query: str) -> str:
        try:
            serpapi_key = os.getenv("SERPAPI_KEY")
            if not serpapi_key:
                logger.warning("SERPAPI_KEY not set. Search unavailable.")
                return "Search unavailable due to missing API key. Please set SERPAPI_KEY in .env."
            if GoogleSearch is None:
                raise ImportError("serpapi not installed")
            params = {
                "engine": "google",
                "q": query,
                "api_key": serpapi_key,
                "num": 8  # Increased results for more info
            }
            search = GoogleSearch(params)
            results = search.get_dict()
            snippets = [result.get("snippet", "") for result in results.get("organic_results", [])]
            return "\n\n".join(snippets) if snippets else "No additional information found via search."
        except Exception as e:
            logger.error(f"SerpAPI search failed: {e}")
            return "Search failed. Please check SERPAPI_KEY and try again."

    def _classify_query(self, question: str) -> tuple[str, bool]:
        question = question.lower().strip()
        is_image_query = any(keyword in question for keyword in ["image", "picture", "scan the image", "video", "link"])
        if "jabalpur smart city" in question:
            return "jabalpur", is_image_query
        if "summary" in question or "summarize" in question:
            return "summary", is_image_query
        elif "important topics" in question or "index" in question or "key topics" in question:
            return "topics", is_image_query
        elif "example" in question or "examples" in question:
            return "examples", is_image_query
        elif any(word in question for word in ["name", "place", "animal", "thing"]):
            return "entities", is_image_query
        elif "question number" in question or "question" in question:
            return "question", is_image_query
        elif any(word in question for word in ["link", "links", "video", "videos", "image", "images"]):
            return "multimedia", is_image_query
        elif "about this file" in question:
            return "file_analysis", is_image_query
        return "general", is_image_query

    def _extract_topics(self, context: str) -> List[str]:
        try:
            if yake is None:
                raise ImportError("YAKE not available for keyword extraction")
            kw_extractor = yake.KeywordExtractor(lan="en", n=3, dedupLim=0.9, top=10)
            keywords = kw_extractor.extract_keywords(context)
            return [kw[0] for kw in keywords]
        except Exception as e:
            logger.error(f"Error extracting topics: {e}")
            return []

    def _extract_examples(self, context: str) -> List[str]:
        examples = []
        sentences = context.split("\n")
        for sentence in sentences:
            if "example" in sentence.lower() or "for instance" in sentence.lower() or "e.g." in sentence.lower():
                examples.append(sentence.strip())
        return examples if examples else ["No examples found in the document."]

    @staticmethod
    def _extract_entities(context: str, query: str) -> List[Dict]:
        if nlp is None:
            return [{"text": "spaCy not available for entity extraction", "label": "N/A"}]
        doc = nlp(context)
        entities = []
        query_lower = query.lower()
        wanted_labels = []
        if "name" in query_lower:
            wanted_labels.append("PERSON")
        if "place" in query_lower:
            wanted_labels.append("GPE")
        if "animal" in query_lower:
            wanted_labels.append("NORP")
        if "thing" in query_lower:
            wanted_labels.extend(["PRODUCT", "OBJECT"])
        for ent in doc.ents:
            if not wanted_labels or ent.label_ in wanted_labels:
                if ent.label_ == "PERSON" and len(ent.text.split()) > 2:
                    continue
                entities.append({"text": ent.text, "label": ent.label_})
        return entities if entities else [{"text": "No relevant entities found", "label": "N/A"}]

    def _extract_multimedia(self, documents: List[LangchainDocument], query: str) -> Dict[str, Any]:
        result = "## 🎥 Multimedia Content\n"
        media = {"images": [], "videos": [], "links": []}
        found_content = False
        for doc in documents:
            metadata = doc.metadata
            source = metadata.get("source", "Unknown")
            if "links" in metadata and "link" in query.lower():
                result += f"### 🔗 Links from {source}\n"
                if metadata["links"]:
                    for link in metadata["links"]:
                        result += f"- [{link}]({link})\n"
                        media["links"].append(link)
                    found_content = True
                else:
                    result += "- No links found.\n"
            if "images" in metadata and "image" in query.lower():
                result += f"### 🖼️ Images from {source}\n"
                if metadata["images"]:
                    for i, img_info in enumerate(metadata["images"]):
                        result += f"- Image {i+1}: {img_info['source']} (Format: {img_info['format']}, Dimensions: {img_info['width']}x{img_info['height']})\n"
                        media["images"].append(img_info)
                    found_content = True
                else:
                    result += "- No images found.\n"
            if "videos" in metadata and "video" in query.lower():
                result += f"### 📹 Videos from {source}\n"
                if metadata["videos"]:
                    for i, video in enumerate(metadata["videos"]):
                        result += f"- Video {i+1}: [View Video]({video})\n"
                        media["videos"].append(video)
                    found_content = True
                else:
                    result += "- No videos found.\n"
        if not found_content:
            result += "**❌ No multimedia content found for the query.**"
        return {"text": result, "media": media}

    def _analyze_images(self, images: List[Dict[str, Any]], question: str = "") -> str:
        result = "## 🖼️ Image Analysis\n"
        question_num = None
        if "question number" in question.lower():
            try:
                question_num = int(re.search(r'\d+', question).group())
            except:
                pass
        found = False
        for i, img_info in enumerate(images):
            try:
                img_bytes = img_info["bytes"].getvalue()
                Image.open(BytesIO(img_bytes))  # Validate
                ocr_text = DocumentProcessor._enhanced_ocr(BytesIO(img_bytes))
                if question_num and f"question {question_num}" in ocr_text.lower():
                    found = True
                    result += f"### 📸 Image for Question {question_num}\n"
                    result += f"- **Source**: {img_info['source']}\n"
                    result += f"- **Dimensions**: {img_info['width']}x{img_info['height']} pixels\n"
                    result += f"- **Format**: {img_info['format']}\n"
                    result += f"- **Extracted Text**: {ocr_text if ocr_text else 'No text detected'}\n\n"
                    return result
                result += f"### 📸 Image {i+1} Details\n"
                result += f"- **Source**: {img_info['source']}\n"
                result += f"- **Dimensions**: {img_info['width']}x{img_info['height']} pixels\n"
                result += f"- **Format**: {img_info['format']}\n"
                result += f"- **Extracted Text**: {ocr_text if ocr_text else 'No text detected'}\n\n"
            except Exception as e:
                logger.error(f"Error analyzing image {i}: {str(e)}")
                result += f"### 📸 Image {i+1}\n- **Error**: Unable to analyze image\n\n"
        if question_num and not found:
            result += f"**❌ No image found for Question {question_num}. Would you like to know more about it?**\n"
        return result

    def _analyze_file(self, documents: List[LangchainDocument], images: List[Dict[str, Any]]) -> str:
        if not documents:
            return "**❌ No file content available for analysis. Please upload a file or URL.**"
        result = "## 📄 File Analysis\n"
        for doc in documents:
            source = doc.metadata.get("source", "Unknown")
            result += f"- **File Name**: {source}\n"
            result += f"- **Type**: {os.path.splitext(source)[1].upper()[1:] if source != 'Unknown' else 'Unknown'}\n"
            content_preview = doc.page_content[:500]
            if "OCR unavailable" in content_preview:
                result += f"- **Content**: {content_preview}\n"
                result += f"- **Details**: This appears to be a scanned document or image-based file. To analyze its content, ensure Tesseract is installed at {TESSERACT_PATH}.\n"
            else:
                result += f"- **Content Preview**: {content_preview}{'...' if len(doc.page_content) > 500 else ''}\n"
                result += f"- **Details**: This file contains text content that can be queried for summaries, entities, or specific questions.\n"
            if images:
                result += f"- **Images**: {len(images)} image(s) extracted. Use 'Show images' to view them.\n"
        result += "\n### 💡 Next Steps\n"
        result += "To extract specific details (e.g., names, dates, or summaries), ask a targeted question like 'Summarize the file' or 'What is the transaction ID?'.\n"
        if "OCR unavailable" in result:
            result += f"To enable OCR for scanned documents, ensure Tesseract is installed at {TESSERACT_PATH} and pytesseract is properly configured.\n"
        return result

    def generate_answer(self, question: str, images: List[Dict[str, Any]] = None, raw_documents: List[LangchainDocument] = None) -> Dict[str, Any]:
        if not self.retriever or not self.llm:
            raise ValueError("Vectorstore and LLM must be initialized")
        cache_key = f"{question}:{''.join(doc.page_content[:100] for doc in raw_documents or [])}"
        if cache_key in llm_cache:
            logger.info("Returning cached LLM response")
            answer = llm_cache[cache_key]
            self.memory.save_context({"input": question}, {"output": answer})
            return {
                "answer": f"<div style='font-weight: bold; font-style: italic; background: linear-gradient(to right, #ffe4e6, #fbcfe8); padding: 20px; border-radius: 12px; max-width: 95%; width: 1400px; margin: auto;'>{answer}</div>",  # Changed style to pink gradient for a fresh look
                "sources": [],
                "media": {"images": [], "videos": [], "links": []}
            }
        self.memory.save_context({"input": question}, {"output": ""})
        query_type, is_image_query = self._classify_query(question)
        docs = self.retriever.invoke(question) if not is_image_query else []
        context = "\n\n".join(doc.page_content for doc in docs) if docs else ""
        if not context and not is_image_query and raw_documents:
            context = "\n\n".join(doc.page_content for doc in raw_documents)[:10000]  # Increased fallback context
        # If context is insufficient, use SerpAPI
        if len(context) < 500 and query_type != "jabalpur":  # Increased threshold
            serp_context = self.search_with_serpapi(question)
            if serp_context:
                context += "\n\n🌐 Additional Information from Web:\n" + serp_context
        try:
            answer = ""
            media = {"images": [], "videos": [], "links": []}
            if query_type == "jabalpur":
                answer = "## 🏙️ Jabalpur Smart City Limited\n"
                answer += "Jabalpur Smart City Limited is an initiative for urban development in Jabalpur. 🌆\n"
                answer += "[Visit Jabalpur Smart City Limited](https://jabalpursmartcity.in/)\n"
                jabalpur_img = fetch_jabalpur_image()
                if jabalpur_img:
                    media["images"].append({"bytes": jabalpur_img, "source": "Jabalpur Logo", "format": "png", "width": 300, "height": 300})
                    answer += "\n![Jabalpur Logo](embedded_image)\n"  # Placeholder for display
            elif query_type == "multimedia" and raw_documents:
                multimedia_result = self._extract_multimedia(raw_documents, question)
                answer = multimedia_result["text"]
                media = multimedia_result["media"]
            elif query_type == "file_analysis" and raw_documents:
                answer = self._analyze_file(raw_documents, images)
            elif is_image_query and images:
                answer = self._analyze_images(images, question)
            else:
                if not context:
                    answer = "**❌ No relevant information found in the file. 💡 Would you like to know more about it? Suggest searching the web or uploading more documents.**"
                    llm_cache[cache_key] = answer
                    self.memory.save_context({"input": question}, {"output": answer})
                    return {
                        "answer": f"<div style='font-weight: bold; font-style: italic; background: linear-gradient(to right, #ffe4e6, #fbcfe8); padding: 20px; border-radius: 12px; max-width: 95%; width: 1400px; margin: auto;'>{answer}</div>",
                        "sources": [],
                        "media": media
                    }
                history = self.memory.load_memory_variables({})["history"]
                history_str = "\n".join([f"User: {msg.content}" if msg.type == "human" else f"Assistant: {msg.content}" for msg in history][-6:])
                template = """
                You are a distinguished professor with expertise in document analysis. Provide a comprehensive, structured response to the query in markdown format, using headings, subheadings, bullet points, tables, emojis for emphasis, and citations where appropriate. Make it engaging and beautiful, with sections like 📋 Overview, 🔑 Key Details, 🌟 Importance, 📍 Location, 📚 Further Reading, and 💡 Suggested Actions. Explain concepts clearly, provide examples, and suggest further reading or actions if information is limited. Include relevant context from the conversation history. If no relevant information is found, state: "❌ No relevant information found. 💡 Suggested action: [suggestion]".

                Conversation History:
                {history}

                Content:
                {context}

                Query: {question}

                Response:
                """
                prompt = ChatPromptTemplate.from_template(template)
                chain = (
                    {"context": lambda x: context, "question": RunnablePassthrough(), "history": lambda x: history_str}
                    | prompt
                    | self.llm
                    | StrOutputParser()
                )
                response = chain.invoke(question)
                answer = response if response.strip() and "no relevant information" not in response.lower() else "**❌ No relevant information found in the file. 💡 Suggested action: Try a different query or upload more documents.**"
            llm_cache[cache_key] = answer
            self.memory.save_context({"input": question}, {"output": answer})
            styled_answer = f"<div style='font-weight: bold; font-style: italic; background: linear-gradient(to right, #ffe4e6, #fbcfe8); padding: 20px; border-radius: 12px; max-width: 95%; width: 1400px; margin: auto;'>{answer}</div>"
            return {
                "answer": styled_answer,
                "sources": [],
                "media": media
            }
        except Exception as e:
            logger.error(f"Error generating answer: {e}")
            answer = f"**❌ Sorry, I couldn't process your request: {str(e)}. Please try again.**"
            llm_cache[cache_key] = answer
            self.memory.save_context({"input": question}, {"output": answer})
            return {
                "answer": f"<div style='font-weight: bold; font-style: italic; background: linear-gradient(to right, #ffe4e6, #fbcfe8); padding: 20px; border-radius: 12px; max-width: 95%; width: 1400px; margin: auto;'>{answer}</div>",
                "sources": [],
                "media": media
            }

# Feedback function with name, email, rating, and feedback text
def send_feedback(name: str, email: str, rating: int, feedback: str):
    try:
        gmail_address = os.getenv("FEEDBACK_EMAIL", "shiveshp449@gmail.com")
        gmail_password = os.getenv("FEEDBACK_PASSWORD", "lyky iugj dvjf crpq")  # App password
        msg = MIMEMultipart("alternative")
        msg['From'] = gmail_address
        msg['To'] = gmail_address
        msg['Subject'] = "DocIntel AI Feedback"
        
        # Professional HTML template for feedback email
        html_body = f"""
        <html>
        <body style="font-family: Arial, sans-serif; color: #333; background-color: #f4f4f4; padding: 20px;">
            <div style="max-width: 600px; margin: auto; background: white; padding: 30px; border-radius: 10px; box-shadow: 0 4px 15px rgba(0,0,0,0.1);">
                <h2 style="color: #4facfe; text-align: center;">DocIntel AI Feedback Received</h2>
                <p style="font-size: 16px; text-align: center;">Thank you for the feedback!</p>
                <hr style="border: 1px solid #eee;">
                <p><strong>Name:</strong> {name}</p>
                <p><strong>Email:</strong> {email}</p>
                <p><strong>Rating:</strong> {'⭐' * rating} ({rating} Stars)</p>
                <p><strong>Feedback:</strong></p>
                <p style="background: #f9f9f9; padding: 15px; border-radius: 5px;">{feedback}</p>
                <hr style="border: 1px solid #eee;">
                <p style="text-align: center; font-size: 14px; color: #888;">DocIntel AI - Unlocking Document Insights</p>
            </div>
        </body>
        </html>
        """
        msg.attach(MIMEText(html_body, 'html'))
        
        server = smtplib.SMTP('smtp.gmail.com', 587)
        server.starttls()
        server.login(gmail_address, gmail_password)
        server.sendmail(gmail_address, gmail_address, msg.as_string())
        server.quit()
        return True
    except Exception as e:
        logger.error(f"Error sending feedback: {e}")
        return False

def send_thank_you_email(name: str, user_email: str):
    try:
        gmail_address = os.getenv("FEEDBACK_EMAIL", "shiveshp449@gmail.com")
        gmail_password = os.getenv("FEEDBACK_PASSWORD", "lyky iugj dvjf crpq")  # App password
        msg = MIMEMultipart("alternative")
        msg['From'] = gmail_address
        msg['To'] = user_email
        msg['Subject'] = "Thank You for Your Feedback on DocIntel AI"
        
        # Professional HTML template for thank-you email
        html_body = f"""
        <html>
        <body style="font-family: Arial, sans-serif; color: #333; background-color: #f4f4f4; padding: 20px;">
            <div style="max-width: 600px; margin: auto; background: white; padding: 30px; border-radius: 10px; box-shadow: 0 4px 15px rgba(0,0,0,0.1);">
                <h2 style="color: #4facfe; text-align: center;">Thank You for Your Feedback!</h2>
                <p style="font-size: 16px; text-align: center;">We appreciate your time and input on DocIntel AI. ⭐</p>
                <hr style="border: 1px solid #eee;">
                <p>Dear {name},</p>
                <p>Your feedback helps us improve our service. We're committed to providing the best document analysis experience.</p>
                <p>If you have any further suggestions, feel free to reply to this email.</p>
                <hr style="border: 1px solid #eee;">
                <p style="text-align: center; font-size: 14px; color: #888;">DocIntel AI - Unlocking Document Insights</p>
            </div>
        </body>
        </html>
        """
        msg.attach(MIMEText(html_body, 'html'))
        
        server = smtplib.SMTP('smtp.gmail.com', 587)
        server.starttls()
        server.login(gmail_address, gmail_password)
        server.sendmail(gmail_address, user_email, msg.as_string())
        server.quit()
        return True
    except Exception as e:
        logger.error(f"Error sending thank-you email: {e}")
        return False

def load_lottie_url(url: str):
    try:
        r = requests.get(url)
        if r.status_code != 200:
            return None
        return r.json()
    except:
        return None

# Main App with Feedback section on main screen and loading message
def main():
    """
    Main function for the DocIntel AI Streamlit app.
    Includes a Feedback section on the main screen with email, rating (star icons), and feedback text inputs.
    Clears the feedback form after submission.
    Sends professional HTML emails for feedback and thank-you.
    Shows 'Analyzing...' during query processing.
    To enable OCR for scanned PDFs/images (e.g., 'NPTEL payment receipt.pdf'):
    1. Install Tesseract:
       - Windows: Download from https://github.com/UB-Mannheim/tesseract/wiki
                 Ensure it is installed at C:\Program Files\Tesseract-OCR
       - Mac: `brew install tesseract`
       - Linux: `sudo apt-get install tesseract-ocr`
    2. Install pytesseract: `pip install pytesseract`
    3. Set TESSERACT_PATH in .env file or environment:
       TESSERACT_PATH=C:\Program Files\Tesseract-OCR\tesseract.exe  # Windows
       TESSERACT_PATH=/usr/bin/tesseract  # Linux/Mac
    """
    st.set_page_config(page_title="DocIntel AI - Advanced Document Analyzer", page_icon="📄", layout="wide")
    
    # Initialize session state variables
    session_vars = {
        "show_welcome": True,
        "chat_history": [],
        "processed_docs": False,
        "vectorstore": None,
        "documents": [],
        "ner_results": None,
        "keyword_results": None,
        "conversion_result": None,
        "extracted_images": [],
        "feedback_name": "",
        "feedback_email": "",
        "feedback_rating": 5,
        "feedback_text": ""
    }
    for key, default in session_vars.items():
        if key not in st.session_state:
            st.session_state[key] = default
    
    # Custom CSS for styling with improved animations and Jarvis animation
    st.markdown("""
    <style>
    .stApp {
        background: linear-gradient(135deg, #0f172a 0%, #1e40af 100%);
        color: white;
    }
    .sidebar .sidebar-content {
        background: rgba(255, 255, 255, 0.95) !important;
        backdrop-filter: blur(12px);
        border-right: 1px solid rgba(255, 255, 255, 0.25);
        padding: 20px;
        border-radius: 10px;
    }
    .stButton>button {
        background: linear-gradient(to right, #3b82f6 0%, #06b6d4 100%);
        color: white;
        border: none;
        border-radius: 30px;
        padding: 12px 32px;
        font-weight: 600;
        box-shadow: 0 4px 15px rgba(59, 130, 246, 0.3);
        transition: all 0.3s ease;
    }
    .stButton>button:hover {
        transform: translateY(-3px) scale(1.05);
        box-shadow: 0 6px 20px rgba(59, 130, 246, 0.4);
        background: linear-gradient(to right, #2563eb 0%, #0e7490 100%);
    }
    .stTextInput>div>div>input, .stTextArea textarea {
        background: rgba(255, 255, 255, 0.95);
        border-radius: 12px;
        padding: 14px;
        border: 1px solid rgba(59, 130, 246, 0.3);
        color: #1e293b;
        font-size: 1.1rem;
        transition: border 0.3s ease;
    }
    .stTextInput>div>div>input:focus, .stTextArea textarea:focus {
        border: 1px solid #3b82f6;
        box-shadow: 0 0 10px rgba(59, 130, 246, 0.3);
    }
    .stSelectbox>div>div>select {
        background: rgba(255, 255, 255, 0.95);
        border-radius: 12px;
        padding: 14px;
        border: 1px solid rgba(59, 130, 246, 0.3);
        color: #1e293b;
        font-size: 1.1rem;
    }
    .stChatMessage {
        border-radius: 20px !important;
        padding: 18px !important;
        margin-bottom: 18px !important;
        box-shadow: 0 4px 15px rgba(0, 0, 0, 0.15);
        animation: fadeInUp 0.5s ease-out;
    }
    @keyframes fadeInUp {
        from { opacity: 0; transform: translateY(20px); }
        to { opacity: 1; transform: translateY(0); }
    }
    .user-message {
        background: linear-gradient(to right, #3b82f6 0%, #06b6d4 100%);
        color: white;
        max-width: 95%;
        width: 1400px;
        margin: auto;
    }
    .assistant-message {
        background: white;
        color: #1e293b;
        font-weight: bold;
        font-style: italic;
        max-width: 95%;
        width: 1400px;
        margin: auto;
    }
    .document-preview {
        background: rgba(255, 255, 255, 0.95);
        border-radius: 15px;
        padding: 25px;
        margin-bottom: 25px;
        box-shadow: 0 4px 15px rgba(0, 0, 0, 0.15);
        color: #1e293b;
        animation: fadeIn 1s ease-in-out;
    }
    .entity-container {
        background: white;
        padding: 25px;
        border-radius: 15px;
        margin-bottom: 25px;
        color: #1e293b;
        animation: fadeIn 1s ease-in-out;
    }
    .entity-table {
        width: 100%;
        border-collapse: separate;
        border-spacing: 0 8px;
        margin-top: 15px;
    }
    .entity-table th, .entity-table td {
        border: none;
        padding: 12px;
        text-align: left;
        transition: background-color 0.3s ease;
    }
    .entity-table th {
        background-color: #3b82f6;
        color: white;
        border-radius: 8px 8px 0 0;
    }
    .entity-table tr {
        background-color: #f8fafc;
        border-radius: 8px;
    }
    .entity-table tr:hover {
        background-color: #e0f7fa;
    }
    .welcome-container {
        animation: fadeIn 1.5s ease-in-out;
    }
    .welcome-title {
        animation: slideIn 1.2s ease-out, pulse 2s infinite;
    }
    .welcome-subtitle {
        animation: slideIn 1.5s ease-out;
    }
    .welcome-text {
        animation: slideIn 1.8s ease-out;
    }
    @keyframes fadeIn {
        from { opacity: 0; }
        to { opacity: 1; }
    }
    @keyframes slideIn {
        from { transform: translateY(30px); opacity: 0; }
        to { transform: translateY(0); opacity: 1; }
    }
    @keyframes pulse {
        0% { transform: scale(1); }
        50% { transform: scale(1.05); }
        100% { transform: scale(1); }
    }
    .feedback-section {
        background: rgba(255, 255, 255, 0.95);
        padding: 20px;
        border-radius: 15px;
        margin-top: 30px;
        box-shadow: 0 4px 15px rgba(0, 0, 0, 0.15);
        animation: fadeIn 1s ease-in-out;
    }
    .thank-you-message {
        background: linear-gradient(to right, #d4fc79, #96e6a1);
        padding: 20px;
        border-radius: 12px;
        text-align: center;
        font-weight: bold;
        color: #333;
        margin-top: 20px;
        animation: fadeInUp 0.5s ease-out;
    }
    /* Advanced Jarvis Animation with orbiting robot */
    #jarvis {
        position: fixed;
        bottom: 20px;
        right: 20px;
        width: 100px;
        height: 100px;
        z-index: 1000;
    }
    #jarvis-circle {
        position: relative;
        width: 100px;
        height: 100px;
        margin: auto;
    }
    .jarvis-outer {
        background-color: rgba(0,0,0,0);
        border: 5px solid rgba(0,183,229,0.9);
        opacity: .9;
        border-radius: 50%;
        box-shadow: 0 0 35px #2187e7;
        width: 100px;
        height: 100px;
        animation: spinPulse 2s infinite ease-in-out;
    }
    .jarvis-inner {
        background-color: rgba(0,0,0,0);
        border: 5px solid rgba(0,183,229,0.9);
        opacity: .9;
        border-left: 5px solid rgba(0,0,0,0);
        border-right: 5px solid rgba(0,0,0,0);
        border-radius: 50%;
        box-shadow: 0 0 15px #2187e7;
        width: 80px;
        height: 80px;
        position: absolute;
        top: 10px;
        left: 10px;
        animation: spinoffPulse 1s infinite linear;
    }
    .jarvis-robot {
        position: absolute;
        top: 50%;
        left: 50%;
        width: 0;
        height: 0;
        animation: orbit 5s linear infinite;
    }
    .jarvis-robot::before {
        content: '🤖';
        font-size: 24px;
        position: relative;
        display: block;
        transform: translate(-50%, -50%);
        text-shadow: 0 0 10px #2187e7;
    }
    @keyframes spinPulse {
        0% { transform:rotate(160deg); opacity:0; box-shadow:0 0 1px #2187e7;}
        50% { transform:rotate(145deg); opacity:1; }
        100% { transform:rotate(-320deg); opacity:0; }
    }
    @keyframes spinoffPulse {
        0% { transform:rotate(0deg); }
        100% { transform:rotate(360deg);  }
    }
    @keyframes orbit {
        0% {
            transform: rotate(0deg) translateX(55px) rotate(0deg);
        }
        100% {
            transform: rotate(360deg) translateX(55px) rotate(-360deg);
        }
    }
    </style>
    """, unsafe_allow_html=True)

    # Jarvis animation HTML
    st.markdown("""
    <div id="jarvis">
        <div id="jarvis-circle">
            <div class="jarvis-outer"></div>
            <div class="jarvis-inner"></div>
            <div class="jarvis-robot"></div>
        </div>
    </div>
    """, unsafe_allow_html=True)

    if st.session_state.show_welcome:
        if st_lottie is not None:
            lottie_welcome = load_lottie_url("https://assets5.lottiefiles.com/packages/lf20_V9t4XZ.json")  # Example AI animation
            if lottie_welcome:
                st_lottie.st_lottie(lottie_welcome, height=200, key="welcome_lottie")
        st.markdown("""
        <div class="welcome-container" style="display: flex; flex-direction: column; align-items: center; justify-content: center; height: 80vh; text-align: center;">
            <h1 class="welcome-title" style="font-size: 4rem; font-weight: 800; color: #ffffff; margin-bottom: 1.2rem; text-shadow: 0 4px 15px rgba(0,0,0,0.3);">DocIntel AI</h1>
            <h3 class="welcome-subtitle" style="font-size: 2rem; font-weight: 500; color: #f0f0f0; margin-bottom: 1.8rem;">Intelligent Document Analysis Powered by RAG</h3>
            <p class="welcome-text" style="font-size: 1.4rem; color: #e0e0e0; max-width: 750px; line-height: 1.7; margin-bottom: 2.5rem;">
                Unlock insights from your documents, images, and web content with advanced AI. Upload files or URLs to summarize, query, and visualize data effortlessly.
            </p>
            <ul style="font-size: 1.2rem; color: #e0e0e0; text-align: left; max-width: 550px; margin-bottom: 2.5rem; list-style-type: none; padding-left: 0;">
                <li>📄 Support for multiple file types including PDFs, images, and more</li>
                <li>🌐 Seamless web content analysis</li>
                <li>💬 Interactive chat for precise queries</li>
                <li>🔍 Advanced entity recognition and visualization</li>
            </ul>
        </div>
        """, unsafe_allow_html=True)
        col1, col2, col3 = st.columns([1,2,1])
        with col2:
            if st.button("Get Started", key="enter-app", help="Click to enter the app", use_container_width=True):
                st.session_state.show_welcome = False
                st.rerun()
        return

    with st.sidebar:
        st.title("DocIntel AI")
        if option_menu is None:
            st.error("streamlit_option_menu not installed")
            selected = "Document Chat"
        else:
            selected = option_menu(
                menu_title=None,
                options=["Document Chat", "File Converter", "Entity Analysis", "Keyword Analysis", "Feedback"],
                icons=["chat-dots", "file-earmark-arrow-up", "graph-up", "key", "star"],
                default_index=0,
                styles={
                    "container": {"background-color": "transparent", "padding": "0"},
                    "nav-link": {"color": "#1e293b", "font-weight": "500", "padding": "12px"},
                    "nav-link-selected": {"background-color": "#3b82f6", "font-weight": "700", "color": "white"}
                }
            )
        if selected == "Document Chat":
            uploaded_files = st.file_uploader(
                "Upload Documents",
                type=["pdf", "docx", "txt", "csv", "pptx", "xlsx", "jpg", "png", "doc", "rtf", "odt", "tex", "tsv", "xls", "ods", "xlsm", "py", "json", "html", "js", "css", "java", "zip", "tar", "gz", "rar", "7z"],
                accept_multiple_files=True,
                help="Upload supported file types for analysis"
            )
            url = st.text_input("Or enter a URL:", placeholder="https://example.com", help="Enter a valid URL for web content analysis")
            if st.button("Process", disabled=not (uploaded_files or (url and url.startswith(('http://', 'https://')))), use_container_width=True):
                with st.spinner("Processing documents..."):
                    try:
                        st.session_state.chat_history = []
                        if st.session_state.vectorstore:
                            st.session_state.vectorstore.memory.clear()
                        file_paths = []
                        if uploaded_files:
                            temp_dir = tempfile.mkdtemp()
                            for uploaded_file in uploaded_files:
                                file_path = os.path.join(temp_dir, uploaded_file.name)
                                with open(file_path, "wb") as f:
                                    f.write(uploaded_file.getbuffer())
                                file_paths.append(file_path)
                        if url and url.startswith(('http://', 'https://')):
                            file_paths.append(url)
                        st.session_state.extracted_images = []
                        for file_path in file_paths:
                            if not file_path.startswith(('http://', 'https://')):
                                images = DocumentProcessor.extract_images(file_path)
                                st.session_state.extracted_images.extend(images)
                            else:
                                web_content = asyncio.run(fetch_url_content(file_path))
                                st.session_state.extracted_images.extend(web_content["images"])
                        rag = RAGPipeline()
                        documents = asyncio.run(DocumentLoader.load_documents(file_paths))
                        if documents:
                            st.session_state.documents = documents
                            rag.create_vectorstore(documents)
                            rag.set_llm()
                            st.session_state.vectorstore = rag
                            st.session_state.processed_docs = True
                            st.success("Documents processed successfully!")
                            with st.spinner("Performing entity analysis..."):
                                st.session_state.ner_results = NERAnalyzer.analyze_documents(documents)
                            with st.spinner("Performing keyword analysis..."):
                                combined_text = "\n\n".join([doc.page_content for doc in documents])
                                keywords = KeywordAnalyzer.extract_keywords(combined_text)
                                st.session_state.keyword_results = {"keywords": keywords, "text": combined_text}
                        else:
                            st.error("No valid content found. Please try again with valid files or URL.")
                    except Exception as e:
                        st.error(f"Processing error: {str(e)}")
                        logger.error(f"Processing error: {str(e)}")
        elif selected == "File Converter":
            conv_file = st.file_uploader(
                "Select file to convert",
                type=["pdf", "docx", "txt", "csv", "pptx", "xlsx", "jpg", "png", "doc", "rtf", "odt", "tex", "tsv", "xls", "ods", "xlsm", "py", "json", "html", "js", "css", "java", "zip", "tar", "gz", "rar", "7z"]
            )
            target_format = st.selectbox(
                "Convert to",
                options=["txt", "json", "md", "html", "xml", "docx", "pdf", "image_to_pdf", "image_to_docx"],
                index=0
            )
            if st.button("Convert", disabled=not conv_file, use_container_width=True):
                with st.spinner("Converting..."):
                    try:
                        temp_dir = tempfile.mkdtemp()
                        file_path = os.path.join(temp_dir, conv_file.name)
                        with open(file_path, "wb") as f:
                            f.write(conv_file.getbuffer())
                        st.session_state.extracted_images = DocumentProcessor.extract_images(file_path)
                        converted = DocumentProcessor.convert_to_text(file_path, target_format)
                        if converted:
                            st.session_state.conversion_result = {
                                "content": converted,
                                "format": target_format,
                                "filename": f"{os.path.splitext(conv_file.name)[0]}.{target_format.split('_')[-1] if 'image_to' in target_format else target_format}"
                            }
                            st.success("Conversion successful!")
                        else:
                            st.error("Conversion failed. Please check the file.")
                    except Exception as e:
                        st.error(f"Conversion error: {str(e)}")
                        logger.error(f"Conversion error: {str(e)}")
        elif selected == "Entity Analysis":
            if st.session_state.ner_results:
                st.success("Entity analysis available")
            else:
                st.info("Process documents first to analyze entities")
        elif selected == "Keyword Analysis":
            if st.session_state.keyword_results:
                st.success("Keyword analysis available")
            else:
                st.info("Process documents first to analyze keywords")

    # Main content area
    if selected == "Document Chat":
        st.header("Document Chat", divider="blue")
        if st.session_state.processed_docs:
            with st.expander("Preview Documents", expanded=False):
                combined_content = "\n\n".join([doc.page_content for doc in st.session_state.documents])
                st.text_area(
                    "Content Preview",
                    value=combined_content[:15000] + ("..." if len(combined_content) > 15000 else ""),
                    height=350,
                    disabled=True
                )
            if st.session_state.extracted_images:
                st.subheader("Extracted Images", divider="blue")
                cols = st.columns(3)
                for i, img_info in enumerate(st.session_state.extracted_images):
                    with cols[i % 3]:
                        try:
                            img_bytes = img_info["bytes"].getvalue()
                            watermarked_img = create_watermarked_image(img_bytes)
                            st.image(watermarked_img, caption=f"Image {i+1} - {img_info['source']}", use_container_width=True)
                            st.download_button(
                                label="Download",
                                data=watermarked_img,
                                file_name=f"image_{i+1}.png",
                                mime="image/png",
                                key=f"dl_{i}",
                                use_container_width=True
                            )
                            st.caption(f"Dimensions: {img_info['width']}x{img_info['height']} | Format: {img_info['format']}")
                        except Exception as e:
                            logger.error(f"Error displaying image {i}: {str(e)}")
                            st.error(f"Image {i+1} unavailable")
            st.subheader("Query Your Documents", divider="blue")
            for msg in st.session_state.chat_history:
                if msg["role"] == "user":
                    with st.chat_message("user"):
                        st.markdown(f'<div class="user-message">{msg["content"]}</div>', unsafe_allow_html=True)
                else:
                    with st.chat_message("assistant"):
                        st.markdown(f'<div class="assistant-message">{msg["content"]}</div>', unsafe_allow_html=True)
                        if "media" in msg and msg["media"] and any(msg["media"].get(key) for key in ["images", "videos", "links"]):
                            for img_info in msg["media"].get("images", []):
                                try:
                                    st.image(create_watermarked_image(img_info["bytes"].getvalue()), caption=f"Image from {img_info['source']}", use_container_width=True)
                                except:
                                    st.write(f"Image from {img_info['source']} unavailable")
                            for video in msg["media"].get("videos", []):
                                try:
                                    st.video(video)
                                except:
                                    st.markdown(f"[Video]({video})")
            if prompt := st.chat_input("Ask a question (e.g., 'Summarize', 'Question 10?', 'Show images', 'Jabalpur Smart City')"):
                with st.chat_message("user"):
                    st.markdown(f'<div class="user-message">{prompt}</div>', unsafe_allow_html=True)
                st.session_state.chat_history.append({"role": "user", "content": prompt})
                with st.spinner("Analyzing..."):
                    if st_lottie is not None:
                        lottie_analyzing = load_lottie_url("https://assets3.lottiefiles.com/packages/lf20_0yfsb3dh.json")  # Analyzing animation
                        if lottie_analyzing:
                            st_lottie.st_lottie(lottie_analyzing, height=100, key="analyzing_lottie")
                    try:
                        response = st.session_state.vectorstore.generate_answer(
                            prompt,
                            st.session_state.extracted_images,
                            st.session_state.documents
                        )
                        with st.chat_message("assistant"):
                            st.markdown(response["answer"], unsafe_allow_html=True)
                            if "media" in response and response["media"].get("images"):
                                for img_info in response["media"]["images"]:
                                    st.image(img_info["bytes"].getvalue(), caption=f"Image from {img_info['source']}", use_container_width=True)
                            if "multimedia" in st.session_state.vectorstore._classify_query(prompt)[0]:
                                for doc in st.session_state.documents:
                                    if "images" in doc.metadata:
                                        for i, img_info in enumerate(doc.metadata["images"]):
                                            try:
                                                st.image(create_watermarked_image(img_info["bytes"].getvalue()), caption=f"Image {i+1} from {doc.metadata['source']}", use_container_width=True)
                                            except:
                                                st.write(f"Image {i+1} unavailable")
                                    if "videos" in doc.metadata:
                                        for i, video in enumerate(doc.metadata["videos"]):
                                            try:
                                                st.video(video)
                                            except:
                                                st.markdown(f"[Video {i+1}]({video})")
                        st.session_state.chat_history.append({
                            "role": "assistant",
                            "content": response["answer"],
                            "sources": response["sources"],
                            "media": response["media"]
                        })
                    except Exception as e:
                        st.error(f"Request failed: {str(e)}")
                        logger.error(f"Answer generation error: {str(e)}")
        else:
            st.info("Upload and process documents or URL to chat.")
    elif selected == "File Converter":
        st.header("File Converter", divider="blue")
        if st.session_state.conversion_result:
            st.subheader("Conversion Result", divider="blue")
            if st.session_state.extracted_images:
                st.subheader("Extracted Images", divider="blue")
                cols = st.columns(3)
                for i, img_info in enumerate(st.session_state.extracted_images):
                    with cols[i % 3]:
                        try:
                            img_bytes = img_info["bytes"].getvalue()
                            watermarked_img = create_watermarked_image(img_bytes)
                            st.image(watermarked_img, caption=f"Image {i+1} - {img_info['source']}", use_container_width=True)
                        except Exception as e:
                            logger.error(f"Error displaying image {i}: {str(e)}")
                            st.error(f"Image {i+1} unavailable")
            content = st.session_state.conversion_result["content"]
            target_format = st.session_state.conversion_result["format"]
            if target_format == "txt":
                st.text_area("Converted Content", value=content, height=450)
                mime = "text/plain"
                data = content.encode('utf-8')
            elif target_format == "json":
                try:
                    json_obj = json.loads(content)
                    st.json(json_obj)
                    mime = "application/json"
                    data = content.encode('utf-8')
                except:
                    st.text_area("Converted Content", value=content, height=450)
                    mime = "text/plain"
                    data = content.encode('utf-8')
            elif target_format == "md":
                st.markdown(content)
                mime = "text/markdown"
                data = content.encode('utf-8')
            elif target_format == "html":
                st.components.v1.html(content, height=450, scrolling=True)
                mime = "text/html"
                data = content.encode('utf-8')
            elif target_format == "xml":
                st.code(content, language="xml")
                mime = "application/xml"
                data = content.encode('utf-8')
            elif target_format in ["docx", "pdf", "image_to_docx", "image_to_pdf"]:
                mime = {
                    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    "pdf": "application/pdf",
                    "image_to_docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                    "image_to_pdf": "application/pdf"
                }[target_format]
                st.write("Download the converted file:")
                data = content.getvalue()
            else:
                mime = "text/plain"
                data = content.encode('utf-8')
            st.download_button(
                label="Download File",
                data=data,
                file_name=st.session_state.conversion_result["filename"],
                mime=mime,
                key="download_converted",
                use_container_width=True
            )
        else:
            st.info("Upload a file and choose format to convert.")
    elif selected == "Entity Analysis":
        st.header("Entity Analysis", divider="blue")
        if st.session_state.ner_results:
            st.markdown(NERAnalyzer.visualize_entities(st.session_state.ner_results["entities"]), unsafe_allow_html=True)
            st.subheader("Entity Relationships", divider="blue")
            graph_img = NERAnalyzer.visualize_graph(st.session_state.ner_results["graph"])
            if graph_img:
                st.image(graph_img, use_container_width=True)
            with st.expander("Raw Text Preview"):
                st.text_area(
                    "Document Text",
                    value=st.session_state.ner_results["raw_text"][:10000] + ("..." if len(st.session_state.ner_results["raw_text"]) > 10000 else ""),
                    height=450,
                    disabled=True
                )
        else:
            st.info("Process documents to perform entity analysis")
    elif selected == "Keyword Analysis":
        st.header("Keyword Analysis", divider="blue")
        if st.session_state.keyword_results:
            keywords = st.session_state.keyword_results["keywords"]
            st.subheader("🔑 Top Keywords", divider="blue")
            for kw, score in keywords:
                st.write(f"🔑 {kw} (Score: {score:.4f})")
            wordcloud_img = KeywordAnalyzer.visualize_keywords(keywords)
            if wordcloud_img:
                st.subheader("☁️ Word Cloud", divider="blue")
                st.image(wordcloud_img, use_container_width=True)
            with st.expander("Raw Text Preview"):
                st.text_area(
                    "Document Text",
                    value=st.session_state.keyword_results["text"][:10000] + ("..." if len(st.session_state.keyword_results["text"]) > 10000 else ""),
                    height=450,
                    disabled=True
                )
        else:
            st.info("Process documents to perform keyword analysis")
    elif selected == "Feedback":
        st.header("Feedback", divider="blue")
        st.markdown('<div class="feedback-section">', unsafe_allow_html=True)
        st.subheader("Provide Feedback")
        st.session_state.feedback_name = st.text_input("Your Name:", value=st.session_state.feedback_name, placeholder="Enter your full name", help="Please provide your name for feedback purposes.")
        st.session_state.feedback_email = st.text_input("Your Email:", value=st.session_state.feedback_email, placeholder="Enter your email address", help="Please provide a valid email address.")
        st.session_state.feedback_rating = st.selectbox(
            "Rate Your Experience:",
            options=[1, 2, 3, 4, 5],
            index=st.session_state.feedback_rating - 1,
            format_func=lambda x: f"{'⭐' * x} {x} Star{'s' if x != 1 else ''}",
            help="Rate your experience from 1 to 5 stars."
        )
        st.session_state.feedback_text = st.text_area("Your Feedback:", value=st.session_state.feedback_text, height=200, help="Share your thoughts or report issues with DocIntel AI.")
        if st.button("Send Feedback", use_container_width=True):
            with st.spinner("Sending feedback..."):
                if not st.session_state.feedback_name.strip():
                    st.error("Please provide your name.")
                elif not st.session_state.feedback_email.strip():
                    st.error("Please provide your email.")
                elif not st.session_state.feedback_text.strip():
                    st.error("Please enter some feedback.")
                else:
                    if send_feedback(st.session_state.feedback_name, st.session_state.feedback_email, st.session_state.feedback_rating, st.session_state.feedback_text):
                        if send_thank_you_email(st.session_state.feedback_name, st.session_state.feedback_email):
                            st.markdown('<div class="thank-you-message">Thank you for your feedback! A confirmation email has been sent to you.</div>', unsafe_allow_html=True)
                        else:
                            st.warning("Feedback sent, but failed to send thank-you email.")
                        # Clear form
                        st.session_state.feedback_name = ""
                        st.session_state.feedback_email = ""
                        st.session_state.feedback_rating = 5
                        st.session_state.feedback_text = ""
                        st.rerun()
                    else:
                        st.error("Failed to send feedback. Please try again.")
        st.markdown('</div>', unsafe_allow_html=True)

if __name__ == "__main__":
    load_dotenv()
    main()